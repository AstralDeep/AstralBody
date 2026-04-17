"""Fixtures for file_tools tests.

We deliberately *generate* test artifacts at runtime (PDFs via reportlab,
DOCX via python-docx, etc.) so the repo doesn't carry binary fixtures and
tests stay deterministic.
"""

from __future__ import annotations

import io
import os
import sys
import time
import uuid
from pathlib import Path
from typing import Tuple

import pytest

# Make `from orchestrator.X import Y` resolve.
_BACKEND = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", ".."))
if _BACKEND not in sys.path:
    sys.path.insert(0, _BACKEND)

# Make the backend/tests/attachments StubDatabase reusable here.
sys.path.insert(0, os.path.abspath(os.path.join(_BACKEND, "tests")))

from attachments.conftest import StubDatabase  # noqa: E402

from orchestrator.attachments import store  # noqa: E402
from orchestrator.attachments.repository import AttachmentRepository  # noqa: E402
from agents.general.file_tools import set_database_for_testing  # noqa: E402


@pytest.fixture
def upload_root(tmp_path: Path, monkeypatch) -> Path:
    monkeypatch.setenv("ATTACHMENT_UPLOAD_ROOT", str(tmp_path))
    return tmp_path


@pytest.fixture
def stub_db():
    db = StubDatabase()
    set_database_for_testing(db)
    yield db
    set_database_for_testing(None)


@pytest.fixture
def repo(stub_db) -> AttachmentRepository:
    return AttachmentRepository(stub_db)


def _persist(repo: AttachmentRepository, *, user_id: str, filename: str,
             category: str, extension: str, content_type: str,
             upload_root: Path, payload: bytes) -> str:
    """Write *payload* to disk under the canonical layout, insert a row, return id."""
    aid = str(uuid.uuid4())
    path, size, sha = store.write(
        user_id=user_id, attachment_id=aid, filename=filename,
        chunks=iter([payload]), max_bytes=10 * 1024 * 1024, root=upload_root,
    )
    repo.insert(
        attachment_id=aid, user_id=user_id, filename=filename,
        content_type=content_type, category=category, extension=extension,
        size_bytes=size, sha256=sha,
        storage_path=str(path.relative_to(upload_root)),
    )
    return aid


# ---------------------------------------------------------------------------
# Fixture builders for each file type
# ---------------------------------------------------------------------------


def make_pdf_with_text(text: str = "Hello PDF world") -> bytes:
    """Build a tiny PDF whose first page contains *text*."""
    from reportlab.pdfgen import canvas  # type: ignore

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(100, 750, text)
    c.showPage()
    c.save()
    return buf.getvalue()


def make_pdf_blank() -> bytes:
    """Build a PDF with no extractable text (single blank page)."""
    from reportlab.pdfgen import canvas  # type: ignore

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.showPage()
    c.save()
    return buf.getvalue()


def make_docx(paragraphs: list[str]) -> bytes:
    import docx  # python-docx

    doc = docx.Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_xlsx(rows: list[list[object]]) -> bytes:
    import openpyxl  # type: ignore

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row in rows:
        ws.append(row)
    wb.create_sheet("Notes").append(["only", "for", "presence"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx(slides: list[tuple[str, str]]) -> bytes:
    from pptx import Presentation  # type: ignore

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]
    title_layout = prs.slide_layouts[0]
    for title, body in slides:
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if body:
            slide.placeholders[1].text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_png(width: int = 32, height: int = 32) -> bytes:
    from PIL import Image  # type: ignore

    img = Image.new("RGB", (width, height), color=(80, 120, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def make_rtf(text: str) -> bytes:
    body = (
        r"{\rtf1\ansi\deff0 {\fonttbl{\f0 Courier;}}"
        + text.replace("\n", r"\par ")
        + "}"
    )
    return body.encode("ascii")


def make_odt(text: str) -> bytes:
    from odf.opendocument import OpenDocumentText  # type: ignore
    from odf import text as odftext  # type: ignore

    doc = OpenDocumentText()
    p = odftext.P(text=text)
    doc.text.addElement(p)
    buf = io.BytesIO()
    doc.write(buf)
    return buf.getvalue()


def make_csv(rows: list[list[object]]) -> bytes:
    import csv

    out = io.StringIO()
    w = csv.writer(out)
    for r in rows:
        w.writerow(r)
    return out.getvalue().encode()


# ---------------------------------------------------------------------------
# Medical imaging fixture builders.
# Each builder imports its library lazily so tests can individually skip when
# the library isn't installed in the current env.
# ---------------------------------------------------------------------------


def make_dicom(
    rows: int = 64,
    cols: int = 64,
    frames: int = 1,
    patient_name: str = "ANON^TEST",
    patient_id: str = "12345",
    modality: str = "CT",
) -> bytes:
    """Build a minimal DICOM file with identifiable PHI tags in the header."""
    import numpy as np  # type: ignore
    import pydicom  # type: ignore
    from pydicom.dataset import FileDataset, FileMetaDataset  # type: ignore
    from pydicom.uid import ExplicitVRLittleEndian, generate_uid  # type: ignore

    file_meta = FileMetaDataset()
    file_meta.MediaStorageSOPClassUID = "1.2.840.10008.5.1.4.1.1.2"  # CT Image Storage
    file_meta.MediaStorageSOPInstanceUID = generate_uid()
    file_meta.TransferSyntaxUID = ExplicitVRLittleEndian
    file_meta.ImplementationClassUID = generate_uid()

    ds = FileDataset(
        "<in-memory>", {}, file_meta=file_meta, preamble=b"\0" * 128,
    )
    ds.PatientName = patient_name
    ds.PatientID = patient_id
    ds.PatientBirthDate = "19700101"
    ds.InstitutionName = "TEST CLINIC"
    ds.ReferringPhysicianName = "DR^SMITH"
    ds.StudyDate = "20250101"
    ds.SeriesDate = "20250101"
    ds.AccessionNumber = "ACC-001"
    ds.StudyInstanceUID = generate_uid()
    ds.SeriesInstanceUID = generate_uid()
    ds.SOPInstanceUID = file_meta.MediaStorageSOPInstanceUID
    ds.SOPClassUID = file_meta.MediaStorageSOPClassUID

    ds.Modality = modality
    ds.Manufacturer = "TestCorp"
    ds.BodyPartExamined = "HEAD"
    ds.StudyDescription = "Test study"
    ds.SeriesDescription = "Test series"
    ds.Rows = rows
    ds.Columns = cols
    ds.SamplesPerPixel = 1
    ds.PhotometricInterpretation = "MONOCHROME2"
    ds.BitsAllocated = 16
    ds.BitsStored = 16
    ds.HighBit = 15
    ds.PixelRepresentation = 0

    if frames > 1:
        ds.NumberOfFrames = frames
        arr = (np.random.rand(frames, rows, cols) * 2048).astype(np.uint16)
    else:
        arr = (np.random.rand(rows, cols) * 2048).astype(np.uint16)
    ds.PixelData = arr.tobytes()

    ds.is_little_endian = True
    ds.is_implicit_VR = False

    buf = io.BytesIO()
    pydicom.dcmwrite(buf, ds, write_like_original=False)
    return buf.getvalue()


def make_nifti(shape: Tuple[int, int, int] = (8, 8, 8), gz: bool = False) -> bytes:
    """Build a tiny NIfTI (.nii) payload from a numpy array."""
    import tempfile
    import numpy as np  # type: ignore
    import nibabel as nib  # type: ignore

    data = (np.random.rand(*shape) * 100).astype(np.int16)
    affine = np.eye(4)
    img = nib.Nifti1Image(data, affine)
    suffix = ".nii.gz" if gz else ".nii"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        nib.save(img, tmp.name)
        tmp_path = tmp.name
    try:
        with open(tmp_path, "rb") as fh:
            return fh.read()
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def make_nrrd(shape: Tuple[int, int, int] = (8, 8, 8)) -> bytes:
    """Build a tiny NRRD volume payload."""
    import tempfile
    import numpy as np  # type: ignore
    import SimpleITK as sitk  # type: ignore

    arr = (np.random.rand(*shape) * 100).astype(np.int16)
    img = sitk.GetImageFromArray(arr)
    with tempfile.NamedTemporaryFile(suffix=".nrrd", delete=False) as tmp:
        sitk.WriteImage(img, tmp.name)
        tmp_path = tmp.name
    try:
        with open(tmp_path, "rb") as fh:
            return fh.read()
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def make_mha(shape: Tuple[int, int, int] = (8, 8, 8)) -> bytes:
    """Build a single-file MetaImage (.mha) payload."""
    import tempfile
    import numpy as np  # type: ignore
    import SimpleITK as sitk  # type: ignore

    arr = (np.random.rand(*shape) * 100).astype(np.int16)
    img = sitk.GetImageFromArray(arr)
    with tempfile.NamedTemporaryFile(suffix=".mha", delete=False) as tmp:
        sitk.WriteImage(img, tmp.name)
        tmp_path = tmp.name
    try:
        with open(tmp_path, "rb") as fh:
            return fh.read()
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def make_ome_tiff(shape: Tuple[int, int, int] = (3, 16, 16)) -> bytes:
    """Build a minimal OME-TIFF payload (CYX, channels=3)."""
    import numpy as np  # type: ignore
    import tifffile  # type: ignore

    arr = (np.random.rand(*shape) * 255).astype(np.uint8)
    buf = io.BytesIO()
    # Use `ome=True` so tifffile emits OME-XML; axes describes the dim order.
    tifffile.imwrite(buf, arr, ome=True, metadata={"axes": "CYX"})
    return buf.getvalue()


def make_tiff(width: int = 32, height: int = 32) -> bytes:
    """Build a plain 2-D TIFF (tifffile without OME metadata)."""
    import numpy as np  # type: ignore
    import tifffile  # type: ignore

    arr = (np.random.rand(height, width, 3) * 255).astype(np.uint8)
    buf = io.BytesIO()
    tifffile.imwrite(buf, arr, photometric="rgb")
    return buf.getvalue()


def make_pyramidal_tiff(size: int = 512, levels: int = 3) -> bytes:
    """Build a tiled pyramidal TIFF OpenSlide can read as a WSI.

    The result is saved to a tempfile and returned as bytes — OpenSlide reads
    from a path, but downstream tests go through ``_persist`` which writes the
    bytes back to disk before opening, so the round-trip is fine.
    """
    import tempfile
    import numpy as np  # type: ignore
    import tifffile  # type: ignore

    # Base level: RGB gradient for visual interest.
    base = np.zeros((size, size, 3), dtype=np.uint8)
    for y in range(size):
        base[y, :, 0] = (y * 255) // max(1, size - 1)
    for x in range(size):
        base[:, x, 1] = (x * 255) // max(1, size - 1)
    base[..., 2] = 128

    with tempfile.NamedTemporaryFile(suffix=".tiff", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        with tifffile.TiffWriter(tmp_path, bigtiff=False) as tif:
            cur = base
            for level in range(levels):
                tif.write(
                    cur,
                    photometric="rgb",
                    tile=(256, 256),
                    compression="zlib",
                    subfiletype=0 if level == 0 else 1,
                )
                # Halve resolution for the next pyramid level.
                cur = cur[::2, ::2, :]
                if cur.shape[0] < 2 or cur.shape[1] < 2:
                    break
        with open(tmp_path, "rb") as fh:
            return fh.read()
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
