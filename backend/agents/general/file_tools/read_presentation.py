"""``read_presentation`` tool: PPTX and ODP."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from agents.general.file_tools import resolve_attachment

logger = logging.getLogger("FileTools.read_presentation")


def _parse_slide_range(spec: Optional[str], total: int) -> List[int]:
    if not spec:
        return list(range(total))
    pages: List[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = int(a) if a else 1
            end = int(b) if b else total
        else:
            start = end = int(part)
        for p in range(max(1, start), min(total, end) + 1):
            pages.append(p - 1)
    return pages


def _read_pptx(path: Path, slide_range: Optional[str]) -> Dict[str, Any]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    indices = _parse_slide_range(slide_range, len(prs.slides))
    out: List[Dict[str, Any]] = []
    for idx in indices:
        slide = prs.slides[idx]
        title = ""
        body_parts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not title and shape == slide.shapes.title:
                    title = line
                elif line:
                    body_parts.append(line)
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes = ""
        out.append({
            "slide_number": idx + 1,
            "title": title,
            "text": "\n".join(body_parts),
            "speaker_notes": notes,
        })
    return {"slide_count": len(prs.slides), "slides": out}


def _read_odp(path: Path, slide_range: Optional[str]) -> Dict[str, Any]:
    from odf.opendocument import load  # type: ignore
    from odf import draw  # type: ignore
    from odf import text as odftext  # type: ignore

    doc = load(str(path))
    pages = doc.getElementsByType(draw.Page)
    indices = _parse_slide_range(slide_range, len(pages))
    out = []
    for idx in indices:
        page = pages[idx]
        ps = page.getElementsByType(odftext.P)
        body = "\n".join(p.firstChild.data if p.firstChild else "" for p in ps)
        out.append({
            "slide_number": idx + 1,
            "title": "",
            "text": body,
            "speaker_notes": "",
        })
    return {"slide_count": len(pages), "slides": out}


def read_presentation(
    attachment_id: str,
    slide_range: Optional[str] = None,
    user_id: Optional[str] = None,
    **_ignored: Any,
) -> Dict[str, Any]:
    """Read a presentation attachment (PPTX or ODP) and return slide text."""
    att, path, err = resolve_attachment(attachment_id, user_id)
    if err is not None:
        return err
    base = {"filename": att.filename}
    try:
        if att.extension == "pptx":
            base.update(_read_pptx(path, slide_range))
        elif att.extension == "odp":
            base.update(_read_odp(path, slide_range))
        else:
            return {"error": {
                "code": "unsupported",
                "message": f"read_presentation does not support .{att.extension}",
            }}
    except Exception as exc:
        logger.exception("presentation parse failed")
        return {"error": {"code": "parse_failed", "message": str(exc)}}
    return base


__all__ = ["read_presentation"]
